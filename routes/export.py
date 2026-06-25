import os
from flask import Blueprint, render_template, session, send_file, jsonify, current_app, request
import io
import pandas as pd
from datetime import datetime

bp = Blueprint('export', __name__)

@bp.route('/')
def export_page():
    return render_template('export.html', active_tab='export')

@bp.route('/api/download/<format>', methods=['GET'])
def download(format):
    dataset_id = session.get('dataset_id')
    if not dataset_id:
        return jsonify({'error': 'No dataset loaded'}), 400
    df, name = current_app.dataset_store.load(dataset_id)
    base = os.path.splitext(name)[0] if name else 'data'
    buf = io.BytesIO()
    if format == 'csv':
        df.to_csv(buf, index=False); buf.seek(0)
        return send_file(buf, mimetype='text/csv', download_name=f'{base}.csv', as_attachment=True)
    elif format == 'excel':
        df.to_excel(buf, index=False); buf.seek(0)
        return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            download_name=f'{base}.xlsx', as_attachment=True)
    elif format == 'json':
        df.to_json(buf, orient='records'); buf.seek(0)
        return send_file(buf, mimetype='application/json', download_name=f'{base}.json', as_attachment=True)
    return jsonify({'error': 'Unsupported format'}), 400

@bp.route('/api/report/pdf', methods=['GET'])
def generate_pdf_report():
    """Generate PDF report with data summary and visualizations"""
    dataset_id = session.get('dataset_id')
    if not dataset_id:
        return jsonify({'error': 'No dataset loaded'}), 400
    
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib import colors
        
        df, name = current_app.dataset_store.load(dataset_id)
        
        # Create PDF buffer
        buf = io.BytesIO()
        
        # Create document
        doc = SimpleDocTemplate(buf, pagesize=letter)
        styles = getSampleStyleSheet()
        elements = []
        
        # Title
        title = Paragraph(f"Data Analysis Report: {name}", styles['Title'])
        elements.append(title)
        elements.append(Spacer(1, 12))
        
        # Generated date
        date_text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        elements.append(Paragraph(date_text, styles['Normal']))
        elements.append(Spacer(1, 20))
        
        # Dataset Overview
        elements.append(Paragraph("Dataset Overview", styles['Heading2']))
        overview_data = [
            ['Metric', 'Value'],
            ['Rows', str(len(df))],
            ['Columns', str(len(df.columns))],
            ['Missing Values', str(int(df.isnull().sum().sum()))],
            ['Duplicates', str(int(df.duplicated().sum()))],
            ['Memory Usage', f"{df.memory_usage(deep=True).sum() / 1024:.2f} KB"]
        ]
        
        table = Table(overview_data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        elements.append(table)
        elements.append(Spacer(1, 20))
        
        # Column Information
        elements.append(Paragraph("Column Information", styles['Heading2']))
        col_data = [['Column', 'Type', 'Non-Null Count', 'Unique Values']]
        for col in df.columns:
            col_data.append([
                col,
                str(df[col].dtype),
                str(int(df[col].notna().sum())),
                str(df[col].nunique())
            ])
        
        col_table = Table(col_data)
        col_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        elements.append(col_table)
        
        # Build PDF
        doc.build(elements)
        buf.seek(0)
        
        return send_file(buf, mimetype='application/pdf', 
                        download_name=f'{os.path.splitext(name)[0]}_report.pdf', 
                        as_attachment=True)
    
    except ImportError:
        return jsonify({'error': 'ReportLab not installed. Run: pip install reportlab'}), 500
    except Exception as e:
        return jsonify({'error': f'PDF generation failed: {str(e)}'}), 500

@bp.route('/api/report/ppt', methods=['GET'])
def generate_ppt_report():
    """Generate PowerPoint presentation with data insights"""
    dataset_id = session.get('dataset_id')
    if not dataset_id:
        return jsonify({'error': 'No dataset loaded'}), 400
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        df, name = current_app.dataset_store.load(dataset_id)
        
        # Create presentation
        prs = Presentation()
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = f"Data Analysis Report"
        subtitle.text = f"Dataset: {name}\nGenerated: {datetime.now().strftime('%Y-%m-%d')}"
        
        # Overview slide
        overview_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = overview_slide.shapes.title
        title.text = "Dataset Overview"
        
        content = overview_slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        p = text_frame.paragraphs[0]
        p.text = f"Total Rows: {len(df)}"
        p.font.size = Pt(18)
        
        for item in [
            f"Total Columns: {len(df.columns)}",
            f"Missing Values: {int(df.isnull().sum().sum())}",
            f"Duplicates: {int(df.duplicated().sum())}",
            f"Memory Usage: {df.memory_usage(deep=True).sum() / 1024:.2f} KB"
        ]:
            p_add = text_frame.add_paragraph()
            p_add.text = item
            p_add.font.size = Pt(16)
        
        # Column details slide
        col_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = col_slide.shapes.title
        title.text = "Column Details"
        
        content = col_slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        p = text_frame.paragraphs[0]
        p.text = "Columns:"
        p.font.size = Pt(18)
        p.font.bold = True
        
        for i, col in enumerate(df.columns[:15]):  # Limit to first 15 columns
            p_add = text_frame.add_paragraph()
            p_add.text = f"• {col} ({df[col].dtype})"
            p_add.font.size = Pt(14)
        
        if len(df.columns) > 15:
            p_add = text_frame.add_paragraph()
            p_add.text = f"... and {len(df.columns) - 15} more columns"
            p_add.font.size = Pt(14)
            p_add.font.italic = True
        
        # Save to buffer
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        
        return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                        download_name=f'{os.path.splitext(name)[0]}_presentation.pptx',
                        as_attachment=True)
    
    except ImportError:
        return jsonify({'error': 'python-pptx not installed. Run: pip install python-pptx'}), 500
    except Exception as e:
        return jsonify({'error': f'PowerPoint generation failed: {str(e)}'}), 500

@bp.route('/api/report/notebook', methods=['GET'])
def generate_notebook():
    """Generate Jupyter notebook with analysis code"""
    dataset_id = session.get('dataset_id')
    if not dataset_id:
        return jsonify({'error': 'No dataset loaded'}), 400
    
    try:
        import nbformat as nbf
        
        df, name = current_app.dataset_store.load(dataset_id)
        
        # Create notebook
        nb = nbf.v4.new_notebook()
        nb['metadata'] = {
            'kernelspec': {
                'display_name': 'Python 3',
                'language': 'python',
                'name': 'python3'
            }
        }
        
        cells = []
        
        # Title cell
        cells.append(nbf.v4.new_markdown_cell(f"# Data Analysis Notebook\n**Dataset:** {name}\n\n**Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"))
        
        # Import cell
        cells.append(nbf.v4.new_code_cell("""import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
%matplotlib inline"""))
        
        # Load data cell
        cells.append(nbf.v4.new_code_cell(f"# Load dataset\ndf = pd.read_csv('{name}')\nprint(f'Dataset shape: {{df.shape}}')\ndf.head()"))
        
        # Basic info cell
        cells.append(nbf.v4.new_code_cell("# Dataset information\ndf.info()\n\n# Statistical summary\ndf.describe()"))
        
        # Missing values cell
        cells.append(nbf.v4.new_code_cell("# Missing values\nmissing = df.isnull().sum()\nmissing[missing > 0].sort_values(ascending=False)"))
        
        # Visualization cell
        cells.append(nbf.v4.new_code_cell("""# Correlation heatmap
plt.figure(figsize=(12, 8))
numeric_df = df.select_dtypes(include=[np.number])
if len(numeric_df.columns) > 1:
    sns.heatmap(numeric_df.corr(), annot=True, cmap='coolwarm', center=0)
    plt.title('Correlation Matrix')
    plt.tight_layout()
    plt.show()"""))
        
        # Distribution cell
        cells.append(nbf.v4.new_code_cell("""# Distribution of numeric columns
numeric_cols = df.select_dtypes(include=[np.number]).columns
for col in numeric_cols[:5]:  # First 5 numeric columns
    plt.figure(figsize=(10, 4))
    plt.hist(df[col].dropna(), bins=30, edgecolor='black')
    plt.title(f'Distribution of {col}')
    plt.xlabel(col)
    plt.ylabel('Frequency')
    plt.show()"""))
        
        nb['cells'] = cells
        
        # Convert to JSON
        buf = io.BytesIO()
        nbf.write(nb, buf)
        buf.seek(0)
        
        return send_file(buf, mimetype='application/x-ipynb+json',
                        download_name=f'{os.path.splitext(name)[0]}_analysis.ipynb',
                        as_attachment=True)
    
    except ImportError:
        return jsonify({'error': 'nbformat not installed. Run: pip install nbformat'}), 500
    except Exception as e:
        return jsonify({'error': f'Notebook generation failed: {str(e)}'}), 500
