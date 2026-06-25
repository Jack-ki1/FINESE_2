import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import tempfile
import os
import zipfile
from datetime import datetime
from typing import List, Dict, Optional
import logging
from io import BytesIO

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# --- Check for required dependencies ---
UTILS_AVAILABLE = True
CONFIG_AVAILABLE = True
PPTX_AVAILABLE = True
REPORTLAB_AVAILABLE = True
XLSXWRITER_AVAILABLE = True

try:
    from utils import (
        calculate_data_health_score,
        generate_recommendation_list
    )
except ImportError as e:
    logger.error(f"Failed to import utils: {e}")
    UTILS_AVAILABLE = False

try:
    from config import BRAND_NAME, SECTION_HEADER_CLASS, SECTION_SUBHEADER_CLASS, EXPORT_DATE_FORMAT
except ImportError as e:
    logger.error(f"Failed to import config: {e}")
    CONFIG_AVAILABLE = False

# --- Import PPTX (only if available) ---
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not installed. PowerPoint export disabled.")

# --- Import ReportLab (only if available) ---
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Table, TableStyle, Spacer, Image
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors
except ImportError:
    REPORTLAB_AVAILABLE = False
    logger.warning("reportlab not installed. PDF export disabled.")

# --- Import xlsxwriter (only if available) ---
try:
    import xlsxwriter
except ImportError:
    XLSXWRITER_AVAILABLE = False
    logger.warning("xlsxwriter not installed. Excel export may be limited.")

# --- Constants ---
BRAND_COLOR = "#0ea5a4"  # Match your theme.py


# =============================
# ✅ HELPER FUNCTION: DEFINED FIRST — USED BY ALL BELOW
# =============================
def _generate_recommendation_list(df: pd.DataFrame, scorecard: Dict) -> List[str]:
    """Shared recommendation list used across PDF, PPTX, and MD."""
    if UTILS_AVAILABLE:
        try:
            return generate_recommendation_list(df, scorecard)
        except Exception:
            return ["⚠️ Recommendation system temporarily unavailable."]
    else:
        return ["⚠️ Recommendation system not available due to missing dependencies."]


# =============================
# ✅ PDF GENERATOR (Uses _generate_recommendation_list)
# =============================
def _generate_reportlab_pdf(df: pd.DataFrame, scorecard: Dict) -> Optional[str]:
    """Generate PDF using ReportLab (no external deps, Windows-safe)."""
    if not REPORTLAB_AVAILABLE:
        logger.warning("ReportLab not installed. PDF export disabled.")
        return None

    # Create temporary directory for images (auto-cleanup on exit)
    temp_dir = tempfile.mkdtemp(prefix="finese_export_")
    temp_files = []  # Track all temp files we create

    try:
        # Build PDF path
        pdf_path = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf").name

        doc = SimpleDocTemplate(pdf_path, pagesize=letter)
        styles = getSampleStyleSheet()
        elements = []

        # Title
        elements.append(Paragraph(f"{BRAND_NAME} Executive Data Report", styles["Title"]))
        elements.append(Spacer(1, 12))

        # Scorecard
        final_score_float = float(scorecard['final_score'])
        badge_map = {
            95: "🏆 Perfect Dataset",
            90: "🥇 Data Master",
            80: "🥈 Clean Data Apprentice",
            70: "🥉 Data Novice",
            60: "⚠️ Needs Attention",
            0: "📉 Critical Issues"
        }
        badge = next((v for k, v in badge_map.items() if final_score_float >= k), badge_map[0])
        elements.append(Paragraph(f"📊 Overall Health Score: <font color='{BRAND_COLOR}'><b>{scorecard['final_score']}/100</b></font> <font color='#000'>({badge})</font>", styles["Heading2"]))
        elements.append(Spacer(1, 12))

        # Summary table
        summary_data = [
            ["Metric", "Value"],
            ["Rows", f"{scorecard['metrics']['total_rows']:,}"],
            ["Columns", str(scorecard['metrics']['total_cols'])],
            ["Missing Cells", f"{scorecard['metrics']['missing_cells']:,} ({scorecard['metrics']['missing_pct']}%)"],
            ["Duplicates", f"{scorecard['metrics']['duplicate_rows']:,}"],
        ]
        summary_table = Table(summary_data)
        summary_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('FONTSIZE', (0, 1), (-1, -1), 8),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        elements.append(summary_table)
        elements.append(Spacer(1, 12))

        # Detailed scores
        elements.append(Paragraph("📈 Health Breakdown", styles["Heading2"]))
        detail_data = [
            ["Dimension", "Score"],
            ["Completeness", f"{scorecard['details']['completeness']}%"],
            ["Consistency", f"{scorecard['details']['consistency']}%"],
            ["Accuracy", f"{scorecard['details']['accuracy']}%"],
            ["Uniqueness", f"{scorecard['details']['uniqueness']}%"],
            ["Timeliness", f"{scorecard['details']['timeliness']}%"],
        ]
        detail_table = Table(detail_data)
        detail_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('FONTSIZE', (0, 1), (-1, -1), 8),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        elements.append(detail_table)
        elements.append(Spacer(1, 12))

        # Recommendations
        elements.append(Paragraph("💡 Key Recommendations", styles["Heading2"]))
        for rec in _generate_recommendation_list(df, scorecard):
            elements.append(Paragraph(f"• {rec}", styles["Normal"]))
        elements.append(Spacer(1, 12))

        # Build PDF
        doc.build(elements)
        return pdf_path

    except Exception as e:
        logger.error(f"Failed to generate PDF: {e}")
        # Clean up any temp files we created
        for f in temp_files:
            try:
                os.unlink(f)
            except:
                pass
        return None


# =============================
# ✅ POWERPOINT GENERATOR
# =============================
def _generate_pptx_report(df: pd.DataFrame, scorecard: Dict) -> Optional[str]:
    """Generate PowerPoint presentation with data insights."""
    if not PPTX_AVAILABLE:
        logger.warning("python-pptx not installed. PPTX export disabled.")
        return None

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = f"{BRAND_NAME} Executive Data Report"
        subtitle.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"

        # Scorecard slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "📊 Data Health Scorecard"

        # Add score details
        left = Inches(1)
        top = Inches(2)
        width = Inches(4)
        height = Inches(0.5)
        
        final_score_float = float(scorecard['final_score'])
        badge_map = {
            95: "🏆 Perfect Dataset",
            90: "🥇 Data Master",
            80: "🥈 Clean Data Apprentice",
            70: "🥉 Data Novice",
            60: "⚠️ Needs Attention",
            0: "📉 Critical Issues"
        }
        badge = next((v for k, v in badge_map.items() if final_score_float >= k), badge_map[0])
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = f"Overall Score: {scorecard['final_score']}/100 ({badge})"
        p.font.bold = True
        p.font.size = Pt(16)

        # Recommendations slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "💡 Key Recommendations"

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(3)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        recommendations = _generate_recommendation_list(df, scorecard)
        for i, rec in enumerate(recommendations[:5]):  # Limit to top 5
            p = tf.add_paragraph()
            p.text = rec
            p.level = 0
            p.font.size = Pt(12)

        # Save to temp file
        temp_pptx = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        temp_pptx.close()
        prs.save(temp_pptx.name)
        return temp_pptx.name

    except Exception as e:
        logger.error(f"Failed to generate PPTX: {e}")
        return None


# =============================
# ✅ MARKDOWN GENERATOR
# =============================
def _generate_markdown_report(df: pd.DataFrame, scorecard: Dict) -> str:
    """Generate Markdown report."""
    try:
        final_score_float = float(scorecard['final_score'])
        badge_map = {
            95: "🏆 Perfect Dataset",
            90: "🥇 Data Master",
            80: "🥈 Clean Data Apprentice",
            70: "🥉 Data Novice",
            60: "⚠️ Needs Attention",
            0: "📉 Critical Issues"
        }
        badge = next((v for k, v in badge_map.items() if final_score_float >= k), badge_map[0])

        md = f"""# {BRAND_NAME} Executive Data Report

## 📊 Overview
- **Generated**: {datetime.now().strftime('%B %d, %Y at %H:%M')}
- **Overall Health Score**: `{scorecard['final_score']}/100` {badge}
- **Rows**: {scorecard['metrics']['total_rows']:,}
- **Columns**: {scorecard['metrics']['total_cols']}
- **Missing Values**: {scorecard['metrics']['missing_cells']:,} ({scorecard['metrics']['missing_pct']}%)
- **Duplicates**: {scorecard['metrics']['duplicate_rows']:,}
- **Latest Date**: {scorecard['metrics']['latest_date'] or 'N/A'}

## 📈 Health Breakdown
| Metric       | Score |
|--------------|-------|
| Completeness | {scorecard['details']['completeness']}% |
| Consistency  | {scorecard['details']['consistency']}% |
| Accuracy     | {scorecard['details']['accuracy']}% |
| Uniqueness   | {scorecard['details']['uniqueness']}% |
| Timeliness   | {scorecard['details']['timeliness']}% |

## 💡 Key Recommendations
{"\n".join(f"- {r}" for r in _generate_recommendation_list(df, scorecard))}

## 📋 Sample Columns (Top 10)
"""

        # Add sample data if available
        try:
            from tabulate import tabulate
            sample_data_str = tabulate(df.head(10), headers='keys', tablefmt='pipe')
        except ImportError:
            # Fallback if tabulate is not available
            sample_data_str = "Sample data not available (requires tabulate package)"

        md += sample_data_str

        md += """

---
> _Generated automatically by FINESE Data Intelligence Platform_
"""

        temp_md = tempfile.NamedTemporaryFile(delete=False, suffix=".md")
        temp_md.write(md.encode('utf-8'))
        temp_md.close()
        return temp_md.name

    except Exception as e:
        logger.error(f"Failed to generate Markdown: {e}")
        # Create a minimal markdown file as fallback
        temp_md = tempfile.NamedTemporaryFile(delete=False, suffix=".md")
        temp_md.write(f"# {BRAND_NAME} Report\n\nReport generation failed: {e}".encode('utf-8'))
        temp_md.close()
        return temp_md.name


# =============================
# ✅ ZIP BUNDLE CREATOR
# =============================
def _create_zip_bundle(pdf_path: Optional[str], pptx_path: Optional[str], markdown_path: str, df) -> BytesIO:
    """Create a ZIP bundle with all available reports."""
    zip_buffer = BytesIO()
    
    with zipfile.ZipFile(zip_buffer, "w") as zip_file:
        # Add PDF if available
        if pdf_path and os.path.exists(pdf_path):
            zip_file.write(pdf_path, f"{BRAND_NAME.lower()}_executive_report.pdf")
        
        # Add PPTX if available
        if pptx_path and os.path.exists(pptx_path):
            zip_file.write(pptx_path, f"{BRAND_NAME.lower()}_executive_deck.pptx")
            
        # Add Markdown if available
        if markdown_path and os.path.exists(markdown_path):
            zip_file.write(markdown_path, f"{BRAND_NAME.lower()}_executive_report.md")
            
        # Add CSV of the data
        csv_buffer = BytesIO()
        # Handle both DataFrame and DatasetContext objects
        if hasattr(df, 'filtered_df'):
            df.filtered_df.to_csv(csv_buffer, index=False)
        else:
            df.to_csv(csv_buffer, index=False)
        zip_file.writestr(f"{BRAND_NAME.lower()}_data_export.csv", csv_buffer.getvalue())
        
    zip_buffer.seek(0)
    return zip_buffer


# =============================
# ✅ MAIN EXPORT FUNCTION
# =============================
def render_export_tab(filtered) -> None:
    """
    Unified Export Dashboard: Generate PDF, PPTX, Markdown, and CSV bundles.
    Uses only reliable, dependency-free methods (ReportLab + python-pptx).
    """
    # Check if critical dependencies are available
    if not UTILS_AVAILABLE or not CONFIG_AVAILABLE:
        st.error("❌ Critical dependencies missing. Please check your installation.")
        return
        
    # Check if it's a DatasetContext object and get the DataFrame
    if hasattr(filtered, 'filtered_df'):
        df_to_check = filtered.filtered_df
    else:
        df_to_check = filtered

    if df_to_check is None or df_to_check.empty:
        st.warning("⚠️ No data available to export.")
        return

    st.markdown(f'<div class="{SECTION_HEADER_CLASS}">📤 Export Executive Reports</div>', unsafe_allow_html=True)
    st.markdown(f'<div class="{SECTION_SUBHEADER_CLASS}">One-click exports for stakeholders — PDF, PPTX, ZIP bundle</div>', unsafe_allow_html=True)

    # --- 1. Basic Downloads ---
    st.markdown("### 📥 Quick Exports")
    col1, col2 = st.columns(2)
    with col1:
        st.download_button(
            label="📥 Download Filtered Data (CSV)",
            data=df_to_check.to_csv(index=False),
            file_name=f"{BRAND_NAME.lower()}_filtered_{datetime.now().strftime('%Y%m%d')}.csv",
            mime="text/csv"
        )
    with col2:
        if XLSXWRITER_AVAILABLE:
            xbuf = BytesIO()
            try:
                with pd.ExcelWriter(xbuf, engine="xlsxwriter") as writer:
                    # Handle both DataFrame and DatasetContext objects
                    if hasattr(filtered, 'filtered_df'):
                        filtered.filtered_df.to_excel(writer, sheet_name="Filtered", index=False)
                    else:
                        filtered.to_excel(writer, sheet_name="Filtered", index=False)
                        
                    if 'work_df' in st.session_state and st.session_state.work_df is not None:
                        # Handle both DataFrame and DatasetContext objects for work_df
                        if hasattr(st.session_state.work_df, 'filtered_df'):
                            st.session_state.work_df.filtered_df.to_excel(writer, sheet_name="Working", index=False)
                        else:
                            st.session_state.work_df.to_excel(writer, sheet_name="Working", index=False)
                st.download_button(
                    label="📥 Download Excel (Filtered + Working)",
                    data=xbuf.getvalue(),
                    file_name=f"{BRAND_NAME.lower()}_data_export_{datetime.now().strftime('%Y%m%d')}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )
            except Exception as e:
                st.warning(f"⚠️ Excel export failed: {e}")
                st.info("Fallback: Use CSV export instead.")
        else:
            st.info("Install `xlsxwriter` for Excel export functionality.")

    # --- 2. Executive Report Bundle (PDF + PPTX + MD) ---
    st.markdown("---")
    st.markdown("### 📄 Executive Report Bundle")

    # Show warnings if dependencies are missing
    if not PPTX_AVAILABLE:
        st.warning("⚠️ `python-pptx` not installed. PowerPoint export disabled.")
    if not REPORTLAB_AVAILABLE:
        st.warning("⚠️ `reportlab` not installed. PDF export disabled.")

    # Fetch scorecard (same logic as quality.py)
    # Handle both DataFrame and DatasetContext for work_df
    if 'work_df' in st.session_state and st.session_state.work_df is not None:
        if hasattr(st.session_state.work_df, 'filtered_df'):
            work_df = st.session_state.work_df.filtered_df
        else:
            work_df = st.session_state.work_df
    else:
        work_df = df_to_check  # fallback to the filtered data
        
    if work_df is None:
        st.warning("⚠️ Working data not available.")
        return
        
    try:
        scorecard = calculate_data_health_score(work_df)
    except Exception as e:
        st.error(f"❌ Failed to calculate data health score: {e}")
        return

    # --- Button: Generate Full Bundle ---
    if st.button("🚀 Generate Full Executive Bundle (PDF + PPTX + Markdown)", type="primary"):
        if not REPORTLAB_AVAILABLE and not PPTX_AVAILABLE:
            st.info("ℹ️ Only Markdown and CSV exports available without additional dependencies.")
            
        with st.spinner("Generating premium executive report... This may take 10–30 seconds."):
            # Generate artifacts
            pdf_path = _generate_reportlab_pdf(work_df, scorecard) if REPORTLAB_AVAILABLE else None
            pptx_path = _generate_pptx_report(work_df, scorecard) if PPTX_AVAILABLE else None
            markdown_path = _generate_markdown_report(work_df, scorecard)

            # Create ZIP bundle
            try:
                zip_buffer = _create_zip_bundle(pdf_path, pptx_path, markdown_path, work_df)
                
                # Offer download
                st.success("✅ Executive Bundle Ready!")
                st.download_button(
                    label="⬇️ Download Full Executive Bundle (.zip)",
                    data=zip_buffer.getvalue(),
                    file_name=f"{BRAND_NAME.lower()}_executive_report_{datetime.now().strftime(EXPORT_DATE_FORMAT)}.zip",
                    mime="application/zip",
                    key="download_executive_bundle"
                )

                # Cleanup temp files
                for path in [pdf_path, pptx_path, markdown_path]:
                    if path and os.path.exists(path):
                        try:
                            os.unlink(path)
                        except Exception:
                            pass
            except Exception as e:
                st.error(f"❌ Failed to create export bundle: {e}")
                st.info("Try individual export options below.")

    # --- Separate PPTX Download Button ---
    st.markdown("---")
    st.markdown("### 🖥️ Download PowerPoint (.pptx) Separately")

    if PPTX_AVAILABLE:
        if st.button("⬇️ Download PowerPoint (.pptx) Only"):
            with st.spinner("Generating PowerPoint report..."):
                try:
                    pptx_path = _generate_pptx_report(work_df, scorecard)
                    if pptx_path:
                        with open(pptx_path, "rb") as f:
                            st.download_button(
                                label="⬇️ Download PPTX",
                                data=f.read(),
                                file_name=f"{BRAND_NAME.lower()}_executive_deck_{datetime.now().strftime(EXPORT_DATE_FORMAT)}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                key="download_pptx_only"
                            )
                        try:
                            os.unlink(pptx_path)
                        except Exception:
                            pass
                    else:
                        st.error("❌ Failed to generate PowerPoint.")
                except Exception as e:
                    st.error(f"❌ Error generating PowerPoint: {e}")
    else:
        st.info("Install `python-pptx` to enable standalone PPTX export.")

    # --- Separate PDF Download Button ---
    st.markdown("---")
    st.markdown("### 📄 Download PDF Report Separately")

    if REPORTLAB_AVAILABLE:
        if st.button("⬇️ Download PDF Report Only"):
            with st.spinner("Generating PDF report..."):
                try:
                    pdf_path = _generate_reportlab_pdf(work_df, scorecard)
                    if pdf_path:
                        with open(pdf_path, "rb") as f:
                            st.download_button(
                                label="⬇️ Download PDF",
                                data=f.read(),
                                file_name=f"{BRAND_NAME.lower()}_executive_report_{datetime.now().strftime(EXPORT_DATE_FORMAT)}.pdf",
                                mime="application/pdf",
                                key="download_pdf_only"
                            )
                        try:
                            os.unlink(pdf_path)
                        except Exception:
                            pass
                    else:
                        st.error("❌ Failed to generate PDF.")
                except Exception as e:
                    st.error(f"❌ Error generating PDF: {e}")
    else:
        st.info("Install `reportlab` to enable PDF export.")


# =============================
# ✅ FETCH QUALITY SCORECARD (Helper)
# =============================
# Note: This function is now in utils.py, but we keep it here for backward compatibility
def _fetch_quality_scorecard(df: pd.DataFrame) -> Dict:
    """Reconstruct the quality scorecard from scratch using same logic as quality.py."""
    return calculate_data_health_score(df)